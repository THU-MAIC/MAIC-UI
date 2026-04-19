"""
PPT Processing API Routes

Handles upload and processing of PPT/PPTX/PDF files for conversion to interactive slides
with AI-generated demo insertions.

Shares the existing Document model and infrastructure with pdf_processing.py.
"""

from fastapi import APIRouter, File, UploadFile, HTTPException, Depends, Form, BackgroundTasks
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from sqlalchemy import or_, func
from typing import Optional, Dict, Any, List
import os
import uuid
import json
import time
import logging
from pathlib import Path
import fitz  # PyMuPDF

from ..core.database import get_db, SessionLocal
from ..models.ppt_document import PPTDocument
from ..models.document import MAX_VERSION_HISTORY
from ..models.user import User
from ..services.ppt_processor import process_ppt_background
from ..services.ai_processor import AIProcessor
from ..core.security import get_current_user

# Configure logging
logger = logging.getLogger(__name__)

# Model to provider mapping (same as pdf_processing.py)
CHINESE_MODELS = ["claude-sonnet-4-6", "claude-opus-4-6", "claude-haiku-4-5-20251001",
                  "glm-4.7", "glm-4.6", "glm-4.6v"]


def get_provider_for_model(model: str) -> str:
    """
    Determine the provider type based on the model name.

    Args:
        model: Model name

    Returns:
        Provider type: "chinese" or "english"
    """
    if model in CHINESE_MODELS:
        return "chinese"
    return "english"


def get_ai_processor(generation_mode: str = "fast", ai_model: Optional[str] = None):
    """
    Get AI processor instance based on configuration.

    Args:
        generation_mode: HTML generation mode ("fast" or "heavy")
        ai_model: Specific model to use. If provided, provider is inferred from model name.

    Returns:
        AIProcessor instance
    """
    # If a specific model is provided, infer provider from model
    if ai_model:
        provider_type = get_provider_for_model(ai_model)
        logger.info(f"Using model {ai_model} with inferred provider: {provider_type}")
    else:
        provider_type = os.getenv('AI_PROVIDER', 'chinese').lower()

    # Set model and API key based on provider type
    if provider_type in ['english', 'gemini', 'openai']:
        if provider_type == 'gemini':
            model = ai_model or os.getenv('GEMINI_MODEL', 'gemini-3-pro-image-preview')
            api_key = os.getenv('GEMINI_API_KEY') or os.getenv('ENGLISH_API_KEY') or os.getenv('MIDDLE_TRANSFER_API_KEY')
        elif provider_type == 'openai':
            model = ai_model or os.getenv('OPENAI_MODEL', 'gpt-4.1')
            api_key = os.getenv('OPENAI_API_KEY') or os.getenv('ENGLISH_API_KEY') or os.getenv('MIDDLE_TRANSFER_API_KEY')
        else:  # english (default)
            model = ai_model or os.getenv('ENGLISH_MODEL', 'gpt-4.1')
            api_key = os.getenv('ENGLISH_API_KEY') or os.getenv('MIDDLE_TRANSFER_API_KEY') or os.getenv('OPENAI_API_KEY') or os.getenv('GEMINI_API_KEY')

        provider_config = {
            'provider': 'english',  # Always use unified English provider
            'model': model,
            'api_key': api_key
        }
    elif provider_type == 'chinese':
        # Unified Chinese provider - auto-detect backend from model name
        if ai_model and ai_model.startswith('claude-'):
            # Anthropic model
            model = ai_model
            api_key = os.getenv('ANTHROPIC_API_KEY')
            base_url = os.getenv('ANTHROPIC_BASE_URL')
        else:
            # Zhipu model (default)
            model = 'glm-4.6v'  # Vision model for slide analysis
            api_key = os.getenv('ZHIPU_API_KEY')
            base_url = None

        text_model = ai_model if ai_model in CHINESE_MODELS else os.getenv('ZHIPU_MODEL', 'glm-4.7')
        provider_config = {
            'provider': 'chinese',
            'model': model,
            'api_key': api_key,
            'base_url': base_url,
            'text_model': text_model
        }
    else:
        raise ValueError(f"Unsupported provider: {provider_type}. Supported providers: english, chinese")

    return AIProcessor(provider_config=provider_config, generation_mode=generation_mode)


router = APIRouter()

# Ensure uploads directory exists
UPLOAD_DIR = Path("uploads/ppt")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)


def get_slide_count(file_path: str, file_type: str) -> int:
    """
    Get the number of slides/pages in a PDF or PPTX file.

    Args:
        file_path: Path to the file
        file_type: Type of file ('pdf' or 'pptx')

    Returns:
        Number of slides/pages
    """
    try:
        if file_type == 'pdf':
            doc = fitz.open(file_path)
            slide_count = len(doc)
            doc.close()
            return slide_count
        elif file_type == 'pptx':
            # For PPTX, convert to PDF first and count pages
            # Or use python-pptx to count slides
            from pptx import Presentation
            prs = Presentation(file_path)
            return len(prs.slides)
        else:
            return 0
    except Exception as e:
        logger.warning(f"Failed to get slide count for {file_path}: {e}")
        return 0


@router.post("/upload")
async def upload_ppt(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    title: str = Form(...),
    file_type: str = Form(...),  # 'pptx' or 'pdf'
    subject: Optional[str] = Form(None),
    grade_level: Optional[int] = Form(None),
    description: Optional[str] = Form(None),
    is_public: bool = Form(False),
    user_preferences: Optional[str] = Form(None),
    auto_process: bool = Form(True),  # Whether to auto-process with default settings
    zhipu_text_model: Optional[str] = Form(None),  # Deprecated: use ai_model instead
    ai_model: Optional[str] = Form(None),  # Model selection (e.g., 'glm-4.7', 'claude-sonnet-4-6')
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Upload and process a PPT/PPTX/PDF file to convert it into interactive slides with demos.

    Form Fields:
        file: Uploaded file (PPTX or PDF)
        title: Title for the document
        file_type: Type of file ('pptx' or 'pdf')
        subject: Optional subject/topic
        grade_level: Optional grade level (K-12)
        description: Optional description
        is_public: Whether to make the document public (default: false)
        user_preferences: Optional JSON string with additional preferences
        auto_process: Whether to auto-process with default settings (default: true)
                      Set to false to configure processing options before analysis
    """
    start_time = time.time()
    logger.info(f"Starting PPT upload process for file: {file.filename}, type: {file_type}")

    try:
        # Validate file type
        if file_type not in ['pptx', 'pdf']:
            raise HTTPException(
                status_code=400,
                detail="Invalid file_type. Must be 'pptx' or 'pdf'"
            )

        # Generate unique filename
        file_extension = os.path.splitext(file.filename)[1]
        unique_filename = f"{uuid.uuid4()}{file_extension}"
        file_path = UPLOAD_DIR / unique_filename

        # Save uploaded file
        file_save_start = time.time()
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        file_save_time = time.time() - file_save_start
        logger.info(f"File saved in {file_save_time:.2f}s - Size: {len(content)/1024/1024:.2f}MB")

        # Parse user preferences if provided
        user_prefs = {}
        if user_preferences:
            try:
                user_prefs = json.loads(user_preferences)
            except json.JSONDecodeError:
                logger.warning("Invalid user_preferences JSON, using empty dict")
                user_prefs = {}

        # Add form fields to preferences
        if grade_level:
            user_prefs["grade_level"] = grade_level
        if subject:
            user_prefs["subject"] = subject
        if description:
            user_prefs["description"] = description

        # Handle model selection (ai_model takes precedence over zhipu_text_model)
        selected_model = ai_model or zhipu_text_model
        if selected_model:
            user_prefs["ai_model"] = selected_model
            # Keep zhipu_text_model for backwards compatibility
            user_prefs["zhipu_text_model"] = selected_model

        # Get slide count before creating document record
        slide_count_start = time.time()
        slide_count = get_slide_count(str(file_path), file_type)
        slide_count_time = time.time() - slide_count_start
        logger.info(f"Slide count retrieved in {slide_count_time:.2f}s - {slide_count} pages")

        # Create document record
        db_start = time.time()

        # If not auto_process, set status to "uploaded" waiting for configuration
        status = "processing" if auto_process else "uploaded"

        # Store ai_model in processing_config for later use
        processing_config = {"config_completed": auto_process}
        if selected_model:
            processing_config["ai_model"] = selected_model
            # Keep zhipu_text_model for backwards compatibility
            processing_config["zhipu_text_model"] = selected_model

        document = PPTDocument(
            title=title,
            original_filename=file.filename,
            file_path=str(file_path),
            file_type=file_type,
            file_size=len(content),
            slide_count=slide_count,
            user_id=current_user.id,
            subject=subject,
            grade_level=grade_level,
            description=description,
            is_public=is_public,
            status=status,
            processing_config=processing_config
        )
        db.add(document)
        db.commit()
        db.refresh(document)
        db_time = time.time() - db_start
        logger.info(f"Database record created in {db_time:.2f}s - Document ID: {document.id}")

        # If auto_process is True, add background task to process the PPT
        if auto_process:
            from ..core.database import SessionLocal
            background_tasks.add_task(
                process_ppt_background,
                document_id=document.id,
                file_path=str(file_path),
                file_type=file_type,
                user_prefs=user_prefs,
                db_session_factory=lambda: SessionLocal()
            )
            message = "PPT uploaded successfully. Processing is happening in the background. Use the /ppt/documents/{id}/status endpoint to check progress."
        else:
            message = "PPT uploaded successfully. Please configure processing options to begin analysis."

        total_time = time.time() - start_time
        logger.info(f"PPT upload endpoint completed in {total_time:.2f}s")

        # Return immediately with document info
        return {
            "id": document.id,
            "title": document.title,
            "original_filename": document.original_filename,
            "file_type": document.file_type,
            "subject": document.subject,
            "grade_level": document.grade_level,
            "status": status,
            "message": message,
            "created_at": document.created_at,
            "needs_configuration": not auto_process
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Upload failed: {e}")
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")


@router.post("/documents/{document_id}/configure")
async def configure_ppt_processing(
    document_id: int,
    background_tasks: BackgroundTasks,
    mode: str = Form(...),  # 'batch' or 'specific_pages'
    batch_size: Optional[int] = Form(None),  # For batch mode
    selected_pages: Optional[str] = Form(None),  # JSON array for specific_pages mode
    use_templates: Optional[bool] = Form(True),  # Whether to use template workflow
    ai_model: Optional[str] = Form(None),  # Model selection (e.g., 'glm-4.7', 'claude-sonnet-4-6')
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Configure PPT processing options and start analysis.

    Form Fields:
        mode: Processing mode - 'batch' or 'specific_pages'
        batch_size: Number of slides per batch (for batch mode, default 5)
        selected_pages: JSON array of page numbers to analyze (for specific_pages mode)
        use_templates: Whether to use template workflow (default: True).
                      - True: Search templates and await user selection
                      - False: Generate demos directly with AI (original workflow)
        ai_model: AI model selection (e.g., 'glm-4.7', 'claude-sonnet-4-6').
                  Provider is auto-detected from model name.
                  - Anthropic models: claude-sonnet-4-6, claude-opus-4-6, claude-haiku-4-5-20251001
                  - Zhipu models: glm-4.7, glm-4.6, glm-4.6v
    """
    try:
        # Get document
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.user_id == current_user.id
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        if document.status not in ["uploaded", "processing"]:
            raise HTTPException(
                status_code=400,
                detail=f"Cannot configure document with status: {document.status}"
            )

        # Validate mode
        if mode not in ['batch', 'specific_pages']:
            raise HTTPException(
                status_code=400,
                detail="Invalid mode. Must be 'batch' or 'specific_pages'"
            )

        # Build processing config
        processing_config = {
            "mode": mode,
            "config_completed": True,
            "use_templates": use_templates if use_templates is not None else True
        }

        # Get ai_model from form or from existing processing_config
        selected_model = ai_model or (document.processing_config or {}).get("ai_model")
        if selected_model:
            processing_config["ai_model"] = selected_model
            # Keep zhipu_text_model for backwards compatibility
            processing_config["zhipu_text_model"] = selected_model

        if mode == 'batch':
            processing_config['batch_size'] = batch_size if batch_size else 5
        else:  # specific_pages
            if not selected_pages:
                raise HTTPException(
                    status_code=400,
                    detail="selected_pages is required for specific_pages mode"
                )
            try:
                pages = json.loads(selected_pages)
                if not isinstance(pages, list):
                    raise ValueError("selected_pages must be an array")
                processing_config['selected_pages'] = pages
            except (json.JSONDecodeError, ValueError) as e:
                raise HTTPException(
                    status_code=400,
                    detail=f"Invalid selected_pages format: {str(e)}"
                )

        # Update document with configuration
        document.processing_config = processing_config
        document.status = "processing"
        db.commit()
        db.refresh(document)

        # Start background processing
        from ..core.database import SessionLocal

        # Build user preferences
        user_prefs = {}
        if document.grade_level:
            user_prefs['grade_level'] = document.grade_level
        if document.subject:
            user_prefs['subject'] = document.subject

        # Add ai_model to user preferences
        if selected_model:
            user_prefs["ai_model"] = selected_model
            user_prefs["zhipu_text_model"] = selected_model  # Backwards compatibility

        background_tasks.add_task(
            process_ppt_background,
            document_id=document.id,
            file_path=document.file_path,
            file_type=document.file_type,
            user_prefs=user_prefs,
            db_session_factory=lambda: SessionLocal()
        )

        logger.info(f"Document {document_id} configured with {mode} mode and processing started")

        return {
            "id": document.id,
            "status": "processing",
            "message": f"已开始处理，使用{mode}模式",
            "processing_config": processing_config
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Configuration failed: {e}")
        raise HTTPException(status_code=500, detail=f"Configuration failed: {str(e)}")


@router.get("/documents")
async def get_ppt_documents(
    skip: int = 0,
    limit: int = 20,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Get list of user's PPT presentation documents.
    """
    try:
        documents = db.query(PPTDocument).filter(
            PPTDocument.user_id == current_user.id
        ).order_by(PPTDocument.created_at.desc()).offset(skip).limit(limit).all()

        return [
            {
                "id": doc.id,
                "title": doc.title,
                "original_filename": doc.original_filename,
                "file_type": doc.file_type,
                "slide_count": doc.slide_count,
                "subject": doc.subject,
                "grade_level": doc.grade_level,
                "status": doc.status,
                "created_at": doc.created_at,
                "updated_at": doc.updated_at,
                "root_document_id": doc.root_document_id,
                "version_number": doc.version_number,
                "is_current": doc.is_current
            }
            for doc in documents
        ]

    except Exception as e:
        logger.error(f"Failed to fetch documents: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch documents: {str(e)}")


@router.get("/documents/{document_id}")
async def get_ppt_document(
    document_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Get specific PPT document details including slide data and demos.
    """
    try:
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.user_id == current_user.id
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        result = {
            "id": document.id,
            "title": document.title,
            "original_filename": document.original_filename,
            "file_type": document.file_type,
            "slide_count": document.slide_count,
            "subject": document.subject,
            "grade_level": document.grade_level,
            "description": document.description,
            "status": document.status,
            "created_at": document.created_at,
            "updated_at": document.updated_at
        }

        # Include slides data if available
        if document.slides_data:
            result["slides"] = document.slides_data.get("slides", [])

        # Include analysis results
        if document.analysis_results:
            result["analysis"] = document.analysis_results

        if document.error_message:
            result["error_message"] = document.error_message

        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch document: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch document: {str(e)}")


@router.get("/documents/{document_id}/slides")
async def get_ppt_slides(
    document_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Get all slides for a PPT document with metadata and demo information.
    """
    try:
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.user_id == current_user.id
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        if not document.slides_data:
            return {
                "document_id": document_id,
                "slides": [],
                "message": "Slides not yet processed"
            }

        return {
            "document_id": document_id,
            "title": document.title,
            "slides": document.slides_data.get("slides", [])
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch slides: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch slides: {str(e)}")


@router.get("/documents/{document_id}/interactive-view")
async def get_interactive_view(
    document_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Get complete interactive presentation with slides and demos interspersed.

    Returns an ordered list that the frontend can render sequentially,
    with demo pages automatically inserted after slides that need them.

    If the requested document is part of a version chain, this endpoint
    automatically returns the current version of that chain.
    """
    try:
        # First, find the requested document
        requested_document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id
        ).first()

        if not requested_document:
            raise HTTPException(status_code=404, detail="Document not found")

        # Check access permissions
        if requested_document.user_id != current_user.id and not requested_document.is_public:
            raise HTTPException(status_code=403, detail="Access denied")

        # If this document is part of a version chain, find the current version
        root_id = requested_document.root_document_id or requested_document.id
        current_document = db.query(PPTDocument).filter(
            PPTDocument.root_document_id == root_id,
            PPTDocument.is_current == 1
        ).first()

        # If no current version found in the chain, use the root itself
        if not current_document:
            current_document = db.query(PPTDocument).filter(
                PPTDocument.id == root_id
            ).first()

        if not current_document:
            raise HTTPException(status_code=404, detail="Current version not found")

        document = current_document

        # Log if we're returning a different version than requested
        if document.id != document_id:
            logger.info(f"Returning current version {document.id} (v{document.version_number}) instead of requested document {document_id}")

        if document.status != "ready":
            raise HTTPException(
                status_code=400,
                detail=f"Document not ready. Current status: {document.status}"
            )

        if not document.slides_data:
            raise HTTPException(status_code=404, detail="Slide data not found")

        slides = document.slides_data.get("slides", [])
        slides_by_number = {slide.get("slide_number"): slide for slide in slides}
        allowed_items = []
        allowed_keys = set()

        for slide in slides:
            slide_number = slide.get("slide_number")
            if slide_number is None:
                continue
            # HTML slides should only be added if explicitly in interactive_order
            if slide.get("is_html_upload"):
                html_key = ("html_slide", int(slide_number))
                allowed_items.append({"type": "html_slide", "slide_number": int(slide_number)})
                allowed_keys.add(html_key)
                continue
            slide_key = ("slide", int(slide_number))
            allowed_items.append({"type": "slide", "slide_number": int(slide_number)})
            allowed_keys.add(slide_key)
            if slide.get("needs_demo") and slide.get("demo_html"):
                demo_key = ("demo", int(slide_number))
                allowed_items.append({"type": "demo", "slide_number": int(slide_number)})
                allowed_keys.add(demo_key)

        requested_order = document.slides_data.get("interactive_order") or []
        ordered_items = []
        seen = set()
        for item in requested_order:
            item_type = item.get("type")
            if item_type == "demo":
                item_type = "demo"
            elif item_type == "html_slide":
                item_type = "html_slide"
            else:
                item_type = "slide"
            slide_number = item.get("slide_number")
            if slide_number is None:
                continue
            item_key = (item_type, int(slide_number))
            if item_key in allowed_keys and item_key not in seen:
                ordered_items.append({"type": item_type, "slide_number": int(slide_number)})
                seen.add(item_key)

        for item in allowed_items:
            item_key = (item["type"], int(item["slide_number"]))
            if item_key not in seen:
                ordered_items.append(item)
                seen.add(item_key)

        interactive_items = []
        for item in ordered_items:
            slide = slides_by_number.get(item["slide_number"])
            if not slide:
                continue
            if item["type"] == "slide":
                interactive_items.append({
                    "type": "slide",
                    "slide_number": slide["slide_number"],
                    "image_path": slide["image_path"],
                    "title": slide["title"],
                    "description": slide["description"]
                })
            elif item["type"] == "demo":
                if slide.get("needs_demo") and slide.get("demo_html"):
                    interactive_items.append({
                        "type": "demo",
                        "slide_number": slide["slide_number"],
                        "html": slide["demo_html"],
                        "reason": slide.get("demo_reason", ""),
                        "demo_type": slide.get("demo_type", "visualization")
                    })
            elif item["type"] == "html_slide":
                if slide.get("is_html_upload") and slide.get("uploaded_html"):
                    interactive_items.append({
                        "type": "html_slide",
                        "slide_number": slide["slide_number"],
                        "html": slide["uploaded_html"],
                        "title": slide.get("title", slide.get("uploaded_filename", "Uploaded HTML")),
                        "description": slide.get("description", "")
                    })

        return {
            "document_id": document_id,
            "title": document.title,
            "subject": document.subject,
            "grade_level": document.grade_level,
            "is_public": bool(document.is_public),
            "total_items": len(interactive_items),
            "items": interactive_items
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch interactive view: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch interactive view: {str(e)}")


@router.get("/documents/{document_id}/status")
async def get_ppt_status(
    document_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Get processing status of a PPT document.
    """
    try:
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.user_id == current_user.id
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        # Calculate progress based on status
        if document.status == "ready":
            progress = 100
        elif document.status == "awaiting_template_selection":
            progress = 60
        elif document.status == "processing":
            progress = 30
        else:
            progress = 0

        message = {
            "ready": "Processing complete",
            "awaiting_template_selection": "Template selection required",
            "processing": "Processing presentation...",
            "uploaded": "Configuration required"
        }.get(document.status, "Unknown status")

        if document.error_message:
            message = f"Error: {document.error_message}"

        response_data = {
            "document_id": document.id,
            "status": document.status,
            "progress": progress,
            "message": message,
            "slide_count": document.slide_count,
            "is_public": bool(document.is_public)
        }

        # Include template options if awaiting selection
        if document.status == "awaiting_template_selection" and document.template_options:
            response_data["template_options"] = document.template_options

        return response_data

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch status: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch status: {str(e)}")


@router.post("/documents/{document_id}/select-templates")
async def select_templates_and_generate(
    document_id: int,
    background_tasks: BackgroundTasks,
    template_selections: Dict[str, str],  # {slide_number: template_id}
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Select templates for slides and generate demos.

    Request body should be a JSON object mapping slide numbers to template IDs:
    {
        "2": "sorting_visualization",
        "5": "binary_search_tree",
        ...
    }

    Slides not included will use AI-only generation (no template).
    """
    import asyncio

    try:
        # Get document
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.user_id == current_user.id
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        if document.status != "awaiting_template_selection":
            raise HTTPException(
                status_code=400,
                detail=f"Cannot select templates for document with status: {document.status}"
            )

        logger.info(f"Processing template selections for document {document_id}: {template_selections}")

        # Start background task to generate demos with selected templates
        background_tasks.add_task(
            _generate_demos_with_templates,
            document_id=document_id,
            template_selections=template_selections,
            db_session_factory=lambda: SessionLocal()
        )

        # Update document status
        document.status = "processing"
        db.commit()
        db.refresh(document)

        return {
            "document_id": document_id,
            "status": "processing",
            "message": f"Generating demos for {len(template_selections)} slides with templates...",
            "selected_count": len(template_selections)
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Template selection failed: {e}")
        raise HTTPException(status_code=500, detail=f"Template selection failed: {str(e)}")


async def _generate_demos_with_templates(
    document_id: int,
    template_selections: Dict[str, str],
    db_session_factory
):
    """
    Background task to generate demos using selected templates.

    Args:
        document_id: ID of the PPTDocument
        template_selections: Map of slide_number to template_id
        db_session_factory: Function to create new DB sessions
    """
    logger.info(f"Starting template-based demo generation for document {document_id}")

    # Create new database session
    db = db_session_factory()

    try:
        from ..models.ppt_document import PPTDocument
        from ..services.ppt_processor import PPTDemoAnalyzer

        # Get document
        document = db.query(PPTDocument).filter(PPTDocument.id == document_id).first()
        if not document:
            logger.error(f"Document {document_id} not found")
            return

        # Retrieve ai_model from processing_config if available
        processing_config = document.processing_config or {}
        selected_model = processing_config.get("ai_model") or processing_config.get("zhipu_text_model")

        # Determine provider based on model
        if selected_model:
            provider_type = get_provider_for_model(selected_model)
            logger.info(f"🔄 Using model {selected_model} with provider {provider_type}")
        else:
            provider_type = None
            logger.info("🔄 No ai_model specified, using default provider")

        # Get AI processor with appropriate model (using local get_ai_processor)
        ai_processor = get_ai_processor(ai_model=selected_model)

        # For Zhipu provider, also set the text_model attribute for compatibility
        if selected_model and hasattr(ai_processor.provider, 'text_model'):
            ai_processor.provider.text_model = selected_model
            logger.info(f"🔄 Set Zhipu text model: {ai_processor.provider.text_model}")

        demo_analyzer = PPTDemoAnalyzer(ai_processor, db_session_factory)

        # Build user preferences
        user_prefs = {}
        if document.grade_level:
            user_prefs['grade_level'] = document.grade_level
        if document.subject:
            user_prefs['subject'] = document.subject

        # Add ai_model to user_prefs for consistency
        if selected_model:
            user_prefs["ai_model"] = selected_model
            user_prefs["zhipu_text_model"] = selected_model  # Backwards compatibility

        # Get slides data
        slides_data = document.slides_data.get("slides", [])
        analysis_results = document.analysis_results or {}

        # Create map of analyzed slides
        analyzed_slides_map = {}
        for slide_info in analysis_results.get('slides', []):
            analyzed_slides_map[slide_info['slide_number']] = slide_info

        # Generate demos for each slide with selected template
        for slide_data in slides_data:
            slide_number = slide_data.get('slide_number')
            slide_number_str = str(slide_number)

            if slide_number_str in template_selections:
                # Use template-based generation
                template_id = template_selections[slide_number_str]
                slide_info = analyzed_slides_map.get(slide_number)

                if slide_info:
                    logger.info(f"Generating demo for slide {slide_number} with template {template_id}")
                    try:
                        demo_html = await demo_analyzer.generate_demo_html_with_template(
                            slide_info=slide_info,
                            user_preferences=user_prefs,
                            template_id=template_id
                        )
                        slide_data['demo_html'] = demo_html
                        slide_data['template_used'] = template_id
                        slide_data['generation_method'] = 'template'
                    except Exception as e:
                        logger.error(f"Template generation failed for slide {slide_number}: {e}")
                        # Fallback to AI-only generation
                        demo_html = await demo_analyzer.generate_demo_html(
                            slide_image_path="",
                            slide_info=slide_info,
                            user_preferences=user_prefs
                        )
                        slide_data['demo_html'] = demo_html
                        slide_data['generation_method'] = 'ai_fallback'

            elif slide_data.get('needs_demo'):
                # No template selected, use AI-only generation
                slide_info = analyzed_slides_map.get(slide_number)
                if slide_info:
                    logger.info(f"Generating demo for slide {slide_number} with AI-only (no template)")
                    demo_html = await demo_analyzer.generate_demo_html(
                        slide_image_path="",
                        slide_info=slide_info,
                        user_preferences=user_prefs
                    )
                    slide_data['demo_html'] = demo_html
                    slide_data['generation_method'] = 'ai'

        # Update document with generated demos
        logger.info(f"=" * 80)
        logger.info(f"SAVING GENERATED DEMOS TO DATABASE")
        logger.info(f"Total slides: {len(slides_data)}")
        for i, slide in enumerate(slides_data):
            if slide.get('demo_html'):
                logger.info(f"  Slide {slide['slide_number']}: has demo_html ({len(slide['demo_html'])} chars)")
            elif slide.get('needs_demo'):
                logger.warning(f"  Slide {slide['slide_number']}: needs_demo=True but NO demo_html!")
        logger.info(f"=" * 80)

        # Create a new dict to ensure SQLAlchemy detects the change
        new_slides_data = {"slides": slides_data}

        logger.info(f"Setting document.slides_data...")
        document.slides_data = new_slides_data
        document.status = "ready"
        document.template_options = None  # Clear template options after use

        # Flag the slides_data field as modified for SQLAlchemy
        from sqlalchemy.orm.attributes import flag_modified
        flag_modified(document, "slides_data")
        logger.info(f"Flagged slides_data as modified")

        logger.info(f"Committing changes to database...")
        db.commit()
        db.refresh(document)
        logger.info(f"✅ Database commit completed")

        # Verify the save
        slide_14_data = document.slides_data.get('slides', [{}])[13] if len(document.slides_data.get('slides', [])) > 13 else {}
        logger.info(f"Verifying save: slide 14 data keys = {list(slide_14_data.keys())}")
        logger.info(f"Verifying save: slide 14 has demo_html = {bool(slide_14_data.get('demo_html'))}")
        if slide_14_data.get('demo_html'):
            logger.info(f"Verifying save: slide 14 demo_html length = {len(slide_14_data.get('demo_html', ''))}")

        logger.info(f"Template-based demo generation completed for document {document_id}")

    except Exception as e:
        logger.error(f"Error generating demos with templates: {e}")
        # Update document with error
        try:
            document = db.query(PPTDocument).filter(PPTDocument.id == document_id).first()
            if document:
                document.status = "error"
                document.error_message = str(e)
                db.commit()
        except Exception as db_error:
            logger.error(f"Failed to update error status: {db_error}")

    finally:
        db.close()


@router.get("/public/documents")
async def get_public_ppt_documents(
    skip: int = 0,
    limit: int = 20,
    db: Session = Depends(get_db)
):
    """
    Get list of all public PPT documents for public browsing.
    """
    try:
        documents = db.query(PPTDocument).filter(
            PPTDocument.is_public == True,
            PPTDocument.status == "ready"
        ).order_by(PPTDocument.created_at.desc()).offset(skip).limit(limit).all()

        return [
            {
                "id": doc.id,
                "title": doc.title,
                "original_filename": doc.original_filename,
                "file_type": doc.file_type,
                "slide_count": doc.slide_count,
                "subject": doc.subject,
                "grade_level": doc.grade_level,
                "description": doc.description,
                "created_at": doc.created_at,
                "updated_at": doc.updated_at,
                "root_document_id": doc.root_document_id,
                "version_number": doc.version_number,
                "is_current": doc.is_current
            }
            for doc in documents
        ]

    except Exception as e:
        logger.error(f"Failed to fetch public documents: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch public documents: {str(e)}")


@router.get("/public/documents/{document_id}")
async def get_public_ppt_document(
    document_id: int,
    db: Session = Depends(get_db)
):
    """
    Get public PPT document details for public access.
    Only works for documents marked as public.
    """
    try:
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.is_public == True
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Public document not found")

        result = {
            "id": document.id,
            "title": document.title,
            "original_filename": document.original_filename,
            "file_type": document.file_type,
            "slide_count": document.slide_count,
            "subject": document.subject,
            "grade_level": document.grade_level,
            "description": document.description,
            "status": document.status,
            "created_at": document.created_at
        }

        if document.slides_data:
            result["slides"] = document.slides_data.get("slides", [])

        if document.analysis_results:
            result["analysis"] = document.analysis_results

        return result

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch public document: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch public document: {str(e)}")


@router.get("/public/documents/{document_id}/status")
async def get_public_ppt_status(
    document_id: int,
    db: Session = Depends(get_db)
):
    """
    Get processing status of a public PPT document.
    """
    try:
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.is_public == True
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Public document not found")

        progress = 100 if document.status == "ready" else (50 if document.status == "processing" else 0)
        message = "Processing complete" if document.status == "ready" else \
                 ("Processing presentation..." if document.status == "processing" else "Error occurred")

        if document.error_message:
            message = f"Error: {document.error_message}"

        return {
            "document_id": document.id,
            "status": document.status,
            "progress": progress,
            "message": message,
            "slide_count": document.slide_count,
            "is_public": True
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch public status: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch status: {str(e)}")


@router.get("/public/documents/{document_id}/interactive-view")
async def get_public_interactive_view(
    document_id: int,
    db: Session = Depends(get_db)
):
    """
    Get interactive view for public PPT document.

    If the requested document is part of a version chain, this endpoint
    automatically returns the current version of that chain.
    """
    try:
        # First, find the requested document
        requested_document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.is_public == True
        ).first()

        if not requested_document:
            raise HTTPException(status_code=404, detail="Public document not found")

        # If this document is part of a version chain, find the current version
        root_id = requested_document.root_document_id or requested_document.id
        current_document = db.query(PPTDocument).filter(
            PPTDocument.root_document_id == root_id,
            PPTDocument.is_current == 1,
            PPTDocument.is_public == True
        ).first()

        # If no current version found in the chain, use the root itself
        if not current_document:
            current_document = db.query(PPTDocument).filter(
                PPTDocument.id == root_id,
                PPTDocument.is_public == True
            ).first()

        if not current_document:
            raise HTTPException(status_code=404, detail="Current version not found")

        document = current_document

        # Log if we're returning a different version than requested
        if document.id != document_id:
            logger.info(f"Returning current version {document.id} (v{document.version_number}) instead of requested document {document_id}")

        if document.status != "ready":
            raise HTTPException(
                status_code=400,
                detail=f"Document not ready. Current status: {document.status}"
            )

        if not document.slides_data:
            raise HTTPException(status_code=404, detail="Slide data not found")

        slides = document.slides_data.get("slides", [])
        slides_by_number = {slide.get("slide_number"): slide for slide in slides}
        allowed_items = []
        allowed_keys = set()

        for slide in slides:
            slide_number = slide.get("slide_number")
            if slide_number is None:
                continue
            # HTML slides should only be added if explicitly in interactive_order
            if slide.get("is_html_upload"):
                html_key = ("html_slide", int(slide_number))
                allowed_items.append({"type": "html_slide", "slide_number": int(slide_number)})
                allowed_keys.add(html_key)
                continue
            slide_key = ("slide", int(slide_number))
            allowed_items.append({"type": "slide", "slide_number": int(slide_number)})
            allowed_keys.add(slide_key)
            if slide.get("needs_demo") and slide.get("demo_html"):
                demo_key = ("demo", int(slide_number))
                allowed_items.append({"type": "demo", "slide_number": int(slide_number)})
                allowed_keys.add(demo_key)

        requested_order = document.slides_data.get("interactive_order") or []
        ordered_items = []
        seen = set()
        for item in requested_order:
            item_type = item.get("type")
            if item_type == "demo":
                item_type = "demo"
            elif item_type == "html_slide":
                item_type = "html_slide"
            else:
                item_type = "slide"
            slide_number = item.get("slide_number")
            if slide_number is None:
                continue
            item_key = (item_type, int(slide_number))
            if item_key in allowed_keys and item_key not in seen:
                ordered_items.append({"type": item_type, "slide_number": int(slide_number)})
                seen.add(item_key)

        for item in allowed_items:
            item_key = (item["type"], int(item["slide_number"]))
            if item_key not in seen:
                ordered_items.append(item)
                seen.add(item_key)

        interactive_items = []
        for item in ordered_items:
            slide = slides_by_number.get(item["slide_number"])
            if not slide:
                continue
            if item["type"] == "slide":
                interactive_items.append({
                    "type": "slide",
                    "slide_number": slide["slide_number"],
                    "image_path": slide["image_path"],
                    "title": slide["title"],
                    "description": slide["description"]
                })
            elif item["type"] == "demo":
                if slide.get("needs_demo") and slide.get("demo_html"):
                    interactive_items.append({
                        "type": "demo",
                        "slide_number": slide["slide_number"],
                        "html": slide["demo_html"],
                        "reason": slide.get("demo_reason", ""),
                        "demo_type": slide.get("demo_type", "visualization")
                    })
            elif item["type"] == "html_slide":
                if slide.get("is_html_upload") and slide.get("uploaded_html"):
                    interactive_items.append({
                        "type": "html_slide",
                        "slide_number": slide["slide_number"],
                        "html": slide["uploaded_html"],
                        "title": slide.get("title", slide.get("uploaded_filename", "Uploaded HTML")),
                        "description": slide.get("description", "")
                    })

        return {
            "document_id": document_id,
            "title": document.title,
            "subject": document.subject,
            "grade_level": document.grade_level,
            "is_public": True,
            "total_items": len(interactive_items),
            "items": interactive_items
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to fetch public interactive view: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to fetch public interactive view: {str(e)}")


@router.get("/documents/{document_id}/thumbnails")
async def get_slide_thumbnails(
    document_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Generate slide thumbnails for configuration page.
    Returns list of slides with thumbnail image data (base64).
    """
    import base64
    from io import BytesIO

    try:
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.user_id == current_user.id
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        file_path = document.file_path
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="File not found")

        thumbnails = []

        if document.file_type == 'pdf':
            # Generate thumbnails from PDF
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)

                # Render page to image (thumbnail size)
                mat = fitz.Matrix(0.5, 0.5)  # Scale down for thumbnail
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("png")
                pix = None

                # Convert to base64
                img_base64 = base64.b64encode(img_data).decode('utf-8')

                thumbnails.append({
                    "page_number": page_num + 1,
                    "thumbnail": f"data:image/png;base64,{img_base64}",
                    "selected": False
                })

            doc.close()

        elif document.file_type == 'pptx':
            # For PPTX, try to convert to PDF using LibreOffice, or fall back to python-pptx
            from pptx import Presentation
            import tempfile

            # Create temporary directory for conversion
            temp_dir = tempfile.mkdtemp()
            # LibreOffice will create PDF with the same name as the input file
            pdf_base_name = os.path.splitext(os.path.basename(file_path))[0]
            pdf_path = os.path.join(temp_dir, f"{pdf_base_name}.pdf")

            try:
                # Try to convert PPTX to PDF using LibreOffice
                import subprocess
                libreoffice_available = True

                # Try using LibreOffice installed in the container
                try:
                    # Check if libreoffice is available
                    result = subprocess.run(["which", "libreoffice"], capture_output=True)
                    if result.returncode == 0:
                        logger.info("Using LibreOffice installed in container")

                        # Set environment for proper font rendering
                        import os as os_module
                        env = os_module.environ.copy()
                        env['LANG'] = 'C.UTF-8'

                        # Run LibreOffice conversion
                        convert_cmd = [
                            "libreoffice", "--headless", "--convert-to", "pdf",
                            "--outdir", temp_dir, file_path
                        ]

                        result = subprocess.run(convert_cmd, capture_output=True, timeout=30, env=env)

                        if result.returncode != 0 or not os.path.exists(pdf_path):
                            libreoffice_available = False
                            logger.warning(f"LibreOffice conversion failed: {result.stderr.decode()}")
                        else:
                            logger.info(f"Successfully converted PPTX to PDF")
                    else:
                        libreoffice_available = False
                        logger.info("LibreOffice not found in container, using fallback")

                except (FileNotFoundError, Exception) as e:
                    logger.info(f"LibreOffice not available: {e}, using fallback")
                    libreoffice_available = False

                if libreoffice_available:
                    # Generate thumbnails from converted PDF
                    doc = fitz.open(pdf_path)
                    for page_num in range(len(doc)):
                        page = doc.load_page(page_num)

                        # Render page to image
                        mat = fitz.Matrix(0.5, 0.5)
                        pix = page.get_pixmap(matrix=mat)
                        img_data = pix.tobytes("png")
                        pix = None

                        # Convert to base64
                        img_base64 = base64.b64encode(img_data).decode('utf-8')

                        thumbnails.append({
                            "page_number": page_num + 1,
                            "thumbnail": f"data:image/png;base64,{img_base64}",
                            "selected": False
                        })

                    doc.close()
                else:
                    # Fallback: just return page numbers without images
                    prs = Presentation(file_path)
                    for i in range(len(prs.slides)):
                        thumbnails.append({
                            "page_number": i + 1,
                            "thumbnail": None,
                            "selected": False
                        })

            finally:
                # Clean up temp directory
                import shutil
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)

        return {
            "document_id": document_id,
            "title": document.title,
            "file_type": document.file_type,
            "slide_count": len(thumbnails),
            "thumbnails": thumbnails
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to generate thumbnails: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate thumbnails: {str(e)}")


@router.delete("/documents/{document_id}")
async def delete_ppt_document(
    document_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Delete a PPT document and all versions in the same version chain.
    """
    try:
        document = db.query(PPTDocument).filter(
            PPTDocument.id == document_id,
            PPTDocument.user_id == current_user.id
        ).first()

        if not document:
            raise HTTPException(status_code=404, detail="Document not found")

        # Delete the whole version chain to keep frontend and DB in sync.
        root_id = document.root_document_id or document.id
        chain_documents = db.query(PPTDocument).filter(
            PPTDocument.user_id == current_user.id,
            or_(PPTDocument.id == root_id, PPTDocument.root_document_id == root_id)
        ).all()

        if not chain_documents:
            chain_documents = [document]

        deleted_document_ids: List[int] = []
        import shutil

        for doc in chain_documents:
            # Delete source file from filesystem
            if doc.file_path and os.path.exists(doc.file_path):
                os.remove(doc.file_path)

            # Delete generated slides directory for each version
            slides_dir = Path(f"uploads/ppt/{doc.id}")
            if slides_dir.exists():
                shutil.rmtree(slides_dir)

            deleted_document_ids.append(doc.id)

        # Delete from database
        for doc in chain_documents:
            db.delete(doc)
        db.commit()

        return {
            "message": "Document deleted successfully",
            "deleted_document_ids": deleted_document_ids,
            "deleted_count": len(deleted_document_ids),
            "deleted_root_document_id": root_id
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to delete document: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to delete document: {str(e)}")
